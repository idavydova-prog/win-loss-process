"""
renewal-attrition-analyzer
---------------------------
Filters an attrition Excel export, analyzes the data, and builds a
formatted PowerPoint slide with hyperlinked account names.

Dependencies: openpyxl, python-pptx, lxml (bundled with python-pptx)

Usage:
    python renewal-attrition-analyzer.py \
        --file "/path/to/Attrition Events.xlsx" \
        --threshold 100000 \
        --product "CC" \
        --period "FY27 Q1" \
        --output-dir "/path/to/output"
"""

import argparse
import os
from collections import defaultdict

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Constants ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1F, 0x4E, 0x79)
LINK_BLUE   = RGBColor(0x00, 0x70, 0xC0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
ROW_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ROW_STRIPE  = RGBColor(0xEB, 0xF3, 0xFB)
ROW_RED     = RGBColor(0xFF, 0xED, 0xED)
ROW_RED2    = RGBColor(0xFF, 0xE0, 0xE0)
DIVIDER     = RGBColor(0xCC, 0xCC, 0xCC)

NSMAP_A     = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NSMAP_R     = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
HYPERLINK_RT = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'

OPTY_BASE   = 'https://org62.lightning.force.com/lightning/r/Opportunity/{}/view'

# Expected column indices (0-based)
COL_ID      = 0
COL_NAME    = 1
COL_MONTH   = 2
COL_FP      = 3   # Full/Partial
COL_REASON  = 4
COL_TYPE    = 6   # Oncycle/Offcycle
COL_OPTY    = 7   # Renewal Opportunity ID
COL_ORG62   = 11
COL_AMOUNT  = 12


# ── Step 1: Load & Filter ──────────────────────────────────────────────────

def load_and_filter(filepath, threshold):
    wb = openpyxl.load_workbook(filepath)
    ws = wb.active
    headers = [cell.value for cell in ws[1]]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        amt = row[COL_AMOUNT]
        if amt is None:
            continue
        try:
            val = float(str(amt).replace('$', '').replace(',', ''))
        except ValueError:
            continue
        if val >= threshold:
            rows.append(row)
    return headers, rows


def save_filtered_excel(headers, rows, output_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Attrition Events'

    hdr_font  = Font(bold=True, color='FFFFFF')
    hdr_fill  = PatternFill(start_color='1F4E79', end_color='1F4E79', fill_type='solid')
    hdr_align = Alignment(horizontal='center', vertical='center', wrap_text=True)

    for ci, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=ci, value=h)
        cell.font  = hdr_font
        cell.fill  = hdr_fill
        cell.alignment = hdr_align

    for ri, row in enumerate(rows, 2):
        for ci, val in enumerate(row, 1):
            cell = ws.cell(row=ri, column=ci, value=val)
            cell.alignment = Alignment(wrap_text=True, vertical='top')

    col_widths = {1:18, 2:30, 3:14, 4:10, 5:28, 6:38, 7:12,
                  8:20, 9:50, 10:50, 11:40, 12:45, 13:16}
    for col, width in col_widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width

    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = f'A1:{get_column_letter(len(headers))}1'
    wb.save(output_path)


# ── Step 2: Analyze ────────────────────────────────────────────────────────

def analyze(rows):
    total      = sum(float(str(r[COL_AMOUNT]).replace('$','').replace(',','')) for r in rows)
    full       = [r for r in rows if str(r[COL_FP]).strip() == 'Full']
    partial    = [r for r in rows if str(r[COL_FP]).strip() == 'Partial']

    oncycle, offcycle, mixed = [], [], []
    for r in rows:
        t = str(r[COL_TYPE]) if r[COL_TYPE] else ''
        if 'Oncycle' in t and 'Offcycle' in t:
            mixed.append(r)
        elif 'Oncycle' in t:
            oncycle.append(r)
        else:
            offcycle.append(r)

    reason_totals = defaultdict(lambda: {'count': 0, 'amount': 0.0})
    for r in rows:
        reason = str(r[COL_REASON]).strip() if r[COL_REASON] else 'Unspecified'
        amt    = float(str(r[COL_AMOUNT]).replace('$','').replace(',',''))
        reason_totals[reason]['count']  += 1
        reason_totals[reason]['amount'] += amt

    sorted_reasons = sorted(reason_totals.items(), key=lambda x: -x[1]['amount'])

    # Build renewal opty URL per account name
    opty_map = {}
    for r in rows:
        name    = str(r[COL_NAME]).strip() if r[COL_NAME] else ''
        opty_id = str(r[COL_OPTY]).strip() if r[COL_OPTY] else ''
        if name and opty_id:
            opty_map[name] = OPTY_BASE.format(opty_id)

    return {
        'total':         total,
        'count':         len(rows),
        'logos_lost':    len(full),
        'logos_retained':len(partial),
        'oncycle':       oncycle,
        'offcycle':      offcycle,
        'mixed':         mixed,
        'reasons':       sorted_reasons,
        'opty_map':      opty_map,
    }


def fmt_m(amount):
    return f'${amount/1_000_000:.1f}M'

def fmt_k(amount):
    return f'${amount/1_000:.0f}K'

def short_name(full_name):
    """Return a short display name for inline narrative use."""
    cutoffs = [' LLC', ' Inc.', ', Inc', ', Ltd', ' Corp', ', LLC',
               ' Limited', ' GmbH', ' B.V.', ' AG']
    name = full_name.strip()
    for c in cutoffs:
        name = name.replace(c, '')
    # Trim long names
    parts = name.split()
    return ' '.join(parts[:3]) if len(parts) > 3 else name


# ── Step 3: Build Slide ────────────────────────────────────────────────────

def _add_hyperlink(run, url, slide):
    rId  = slide.part.relate_to(url, HYPERLINK_RT, is_external=True)
    rPr  = run._r.get_or_add_rPr()
    hlink = etree.SubElement(rPr, f'{{{NSMAP_A}}}hlinkClick')
    hlink.set(f'{{{NSMAP_R}}}id', rId)


def _add_rect(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def _plain_run(p, text, size, color=BLACK, bold=False):
    run = p.add_run()
    run.text  = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return run


def _linked_run(p, text, url, slide, size=10):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = LINK_BLUE
    run.font.underline = True
    _add_hyperlink(run, url, slide)
    return run


def _add_paragraph_mixed(tf, segments, size, color, bold, slide, is_first=False):
    """
    segments: list of (text, url_or_None)
    Adds a paragraph mixing plain and hyperlinked runs.
    """
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.space_after = Pt(3)
    for text, url in segments:
        if not text:
            continue
        if url:
            _linked_run(p, text, url, slide, size)
        else:
            _plain_run(p, text, size, color, bold)
    return p


def _parse_and_link(text, opty_map, slide, tf, size, color, bold, is_first=False):
    """Split a line of text into plain/linked segments based on opty_map keys."""
    names = sorted(opty_map.keys(), key=len, reverse=True)
    segments = []
    remaining = text
    while remaining:
        matched = False
        for name in names:
            idx = remaining.find(name)
            if idx == 0:
                segments.append((name, opty_map[name]))
                remaining = remaining[len(name):]
                matched = True
                break
            elif idx > 0:
                segments.append((remaining[:idx], None))
                remaining = remaining[idx:]
                matched = True
                break
        if not matched:
            segments.append((remaining, None))
            break
    _add_paragraph_mixed(tf, segments, size, color, bold, slide, is_first)


def _cycle_line(tf, label, accounts, opty_map, slide, is_first=False):
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.space_after = Pt(2)
    _plain_run(p, label, 10, BLACK)
    for i, r in enumerate(accounts):
        name = str(r[COL_NAME]).strip()
        url  = opty_map.get(name)
        display = short_name(name)
        if url:
            _linked_run(p, display, url, slide, 10)
        else:
            _plain_run(p, display, 10, BLACK)
        if i < len(accounts) - 1:
            _plain_run(p, ', ', 10, BLACK)


def build_slide(rows, analysis, product, period, output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BLUE

    opty_map = analysis['opty_map']

    # ── Title ──
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.65))
    p  = tb.text_frame.paragraphs[0]
    _plain_run(p, f'{period} {product} Attrition >$100K Details', 26, NAVY, bold=True)

    # ── Table ──
    sorted_rows = sorted(rows, key=lambda r: -float(str(r[COL_AMOUNT]).replace('$','').replace(',','')))
    col_widths  = [Inches(3.8), Inches(1.5), Inches(1.7)]
    num_rows    = 1 + len(sorted_rows)
    tbl = slide.shapes.add_table(
        num_rows, 3, Inches(0.3), Inches(0.9), sum(col_widths), Inches(6.3)
    ).table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for ci, h in enumerate(['Customer', 'Full/Partial', 'Attrition Amount']):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p  = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold  = True
        run.font.size  = Pt(10)
        run.font.color.rgb = WHITE

    aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT]
    for ri, row in enumerate(sorted_rows):
        name  = str(row[COL_NAME]).strip()
        ftype = str(row[COL_FP]).strip()
        amt   = float(str(row[COL_AMOUNT]).replace('$','').replace(',',''))
        amt_s = f'${amt:,.0f}'
        if ftype == 'Full':
            bg_color = ROW_RED if ri % 2 == 0 else ROW_RED2
        else:
            bg_color = ROW_WHITE if ri % 2 == 0 else ROW_STRIPE

        for ci, val in enumerate([name, ftype, amt_s]):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            p = cell.text_frame.paragraphs[0]
            p.alignment = aligns[ci]
            if ci == 0 and name in opty_map:
                _linked_run(p, val, opty_map[name], slide, 9)
            else:
                _plain_run(p, val, 9, BLACK)

    for ri in range(num_rows):
        tbl.rows[ri].height = Inches(0.32)

    # ── Insights Panel ──
    IL = Inches(7.7)
    IT = Inches(0.9)
    IW = Inches(5.35)
    IH = Inches(6.3)
    PL = IL + Inches(0.2)
    PW = IW - Inches(0.4)

    _add_rect(slide, IL, IT, IW, IH, WHITE, DIVIDER)
    _add_rect(slide, IL, IT, IW, Inches(0.44), NAVY)

    hdr = slide.shapes.add_textbox(IL, IT, IW, Inches(0.44))
    p   = hdr.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _plain_run(p, 'Key Insights', 13, WHITE, bold=True)

    def divider(top):
        _add_rect(slide, PL, top, PW, Inches(0.02), DIVIDER)

    # Block 1 — Volume
    y = IT + Inches(0.54)
    tb1 = slide.shapes.add_textbox(PL, y, PW, Inches(1.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.space_after = Pt(3)
    _plain_run(p, f'{analysis["count"]} Customers >$100K attrition', 12, BLACK, bold=True)

    p = tf1.add_paragraph(); p.space_after = Pt(3)
    _plain_run(p, f'{analysis["logos_lost"]} Logos Lost  |  {analysis["logos_retained"]} Logos Retained (Partial)', 11, BLACK)

    p = tf1.add_paragraph(); p.space_after = Pt(3)
    _plain_run(p, f'{fmt_m(analysis["total"])} Total Attrition', 11, BLACK)

    _cycle_line(tf1, f'{len(analysis["oncycle"])} Oncycle: ',  analysis['oncycle'],  opty_map, slide)
    _cycle_line(tf1, f'{len(analysis["offcycle"])} Offcycle: ', analysis['offcycle'], opty_map, slide)
    if analysis['mixed']:
        _cycle_line(tf1, f'{len(analysis["mixed"])} Mixed: ', analysis['mixed'], opty_map, slide)

    divider(IT + Inches(2.5))

    # Block 2 — Top Drivers
    tb2 = slide.shapes.add_textbox(PL, IT + Inches(2.6), PW, Inches(1.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]; p.space_after = Pt(3)
    _plain_run(p, 'Top Attrition Drivers:', 11, BLACK, bold=True)

    for reason, data in analysis['reasons'][:7]:
        p = tf2.add_paragraph(); p.space_after = Pt(2)
        label = f'{reason}: {data["count"]} customer{"s" if data["count"] > 1 else ""}  ({fmt_m(data["amount"])})'
        _plain_run(p, label, 10, BLACK)

    divider(IT + Inches(4.25))

    # Block 3 — Competitive Narrative
    # Build short-name opty_map for narrative matching
    short_opty = {short_name(k): v for k, v in opty_map.items()}

    tb3 = slide.shapes.add_textbox(PL, IT + Inches(4.35), PW, Inches(1.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    # Detect Shopify accounts from reason/notes
    shopify_names = [short_name(str(r[COL_NAME])) for r in rows
                     if 'Shopify' in str(r[COL_REASON]) + str(r[8]) + str(r[9])
                     or 'Competitive' in str(r[COL_REASON])]
    bd_names = [short_name(str(r[COL_NAME])) for r in rows
                if 'Bad Debt' in str(r[COL_REASON]) or 'Bankruptcy' in str(r[COL_REASON])]
    oversold_names = [short_name(str(r[COL_NAME])) for r in rows
                      if 'Oversold' in str(r[COL_REASON]) + str(r[5])]

    def narrative_line(tf, prefix, names, is_first=False):
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.space_after = Pt(3)
        _plain_run(p, prefix, 11, NAVY)
        for i, name in enumerate(names):
            url = short_opty.get(name)
            if url:
                _linked_run(p, name, url, slide, 11)
            else:
                _plain_run(p, name, 11, NAVY)
            if i < len(names) - 1:
                _plain_run(p, ', ', 11, NAVY)

    if shopify_names:
        narrative_line(tf3, 'Shopify competitive losses: ', shopify_names, is_first=True)
    if oversold_names:
        narrative_line(tf3, 'Oversold/GMV mismatch: ', oversold_names, is_first=(not shopify_names))
    if bd_names:
        narrative_line(tf3, 'Bad Debt / Bankruptcy: ', bd_names, is_first=(not shopify_names and not oversold_names))

    p = tf3.add_paragraph(); p.space_after = Pt(3)
    _plain_run(p, 'Competitors: Shopify, homegrown solutions', 11, NAVY)

    prs.save(output_path)


# ── CLI Entry Point ────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='Renewal Attrition Analyzer')
    parser.add_argument('--file',      required=True, help='Path to source attrition Excel file')
    parser.add_argument('--threshold', type=float, default=100000, help='Minimum attrition amount (default 100000)')
    parser.add_argument('--product',   default='CC',      help='Product name for slide title (default: CC)')
    parser.add_argument('--period',    default='FY27 Q1', help='Fiscal period label (default: FY27 Q1)')
    parser.add_argument('--output-dir', default=None,     help='Output folder (default: same as input file)')
    args = parser.parse_args()

    out_dir = args.output_dir or os.path.dirname(args.file)
    prefix  = f'{args.period.replace(" ", "_")}_{args.product}'

    print(f'Loading {args.file}...')
    headers, rows = load_and_filter(args.file, args.threshold)
    print(f'  {len(rows)} rows with attrition >= ${args.threshold:,.0f}')

    excel_out = os.path.join(out_dir, f'{prefix}_Attrition_Events.xlsx')
    save_filtered_excel(headers, rows, excel_out)
    print(f'  Saved filtered data: {excel_out}')

    print('Analyzing...')
    analysis = analyze(rows)
    print(f'  Total: {fmt_m(analysis["total"])}  |  Logos lost: {analysis["logos_lost"]}  |  Retained: {analysis["logos_retained"]}')

    print('Building slide...')
    slide_out = os.path.join(out_dir, f'{prefix}_Attrition_Slide.pptx')
    build_slide(rows, analysis, args.product, args.period, slide_out)
    print(f'  Saved slide: {slide_out}')
    print('Done.')


if __name__ == '__main__':
    main()
